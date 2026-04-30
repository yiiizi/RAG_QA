"""
Multi-format document loader.

Supports 12+ file types:
    PDF (.pdf), Word (.docx), Plain text (.txt), Markdown (.md),
    HTML (.html/.htm), CSV (.csv), Excel (.xlsx), PowerPoint (.pptx),
    JSON (.json), EPUB (.epub), Images (.png/.jpg via OCR),
    Code files (.py/.java/.go/.js/.ts/.cpp/.c)

Uses LlamaIndex SimpleDirectoryReader as the backbone, with custom
supplementary loaders for image OCR and advanced PDF table extraction.
"""

from __future__ import annotations

import logging
import tempfile
from pathlib import Path
from typing import Optional

from config.settings import settings

logger = logging.getLogger(__name__)


def _get_Document():
    """Lazy import to avoid requiring llama-index at server startup."""
    try:
        from llama_index.core import Document
        return Document
    except ImportError as e:
        raise ImportError(
            "llama-index not installed. Run: pip install llama-index"
        ) from e


# ── Per-format loaders ──────────────────────────────────────────────

def _load_pdf(file_path: Path) -> str:
    """Load PDF with PyMuPDF (fast) + fallback to PDFPlumber for tables."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(file_path))
        text_parts: list[str] = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"PyMuPDF failed for {file_path}: {e} — trying PDFPlumber")
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            return "\n\n".join(page.extract_text() or "" for page in pdf.pages)


def _load_docx(file_path: Path) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(file_path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _load_txt(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="replace")


def _load_markdown(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="replace")


def _load_html(file_path: Path) -> str:
    from bs4 import BeautifulSoup
    html = file_path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text("\n", strip=True)


def _load_csv(file_path: Path) -> str:
    import pandas as pd
    df = pd.read_csv(str(file_path))
    return df.to_string(index=False)


def _load_excel(file_path: Path) -> str:
    import pandas as pd
    dfs = pd.read_excel(str(file_path), sheet_name=None)
    parts: list[str] = []
    for sheet_name, df in dfs.items():
        parts.append(f"## Sheet: {sheet_name}\n{df.to_string(index=False)}")
    return "\n\n".join(parts)


def _load_pptx(file_path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(file_path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            parts.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _load_json(file_path: Path) -> str:
    import json
    data = json.loads(file_path.read_text(encoding="utf-8", errors="replace"))
    # Flatten nested JSON into readable text
    def _flatten(obj, depth: int = 0) -> list[str]:
        lines: list[str] = []
        prefix = "  " * depth
        if isinstance(obj, dict):
            for k, v in obj.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}- {k}:")
                    lines.extend(_flatten(v, depth + 1))
                else:
                    lines.append(f"{prefix}- {k}: {v}")
        elif isinstance(obj, list):
            for item in obj:
                lines.extend(_flatten(item, depth))
        else:
            lines.append(f"{prefix}{obj}")
        return lines
    return "\n".join(_flatten(data))


def _load_epub(file_path: Path) -> str:
    from bs4 import BeautifulSoup
    from ebooklib import epub, ITEM_DOCUMENT

    book = epub.read_epub(str(file_path))
    parts: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        text = soup.get_text("\n", strip=True)
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _load_image_ocr(file_path: Path) -> str:
    """OCR recognition via PaddleOCR."""
    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(lang="ch", use_angle_cls=True, show_log=False)
        result = ocr.ocr(str(file_path), cls=True)
        lines: list[str] = []
        if result and result[0]:
            for line_info in result[0]:
                if line_info and len(line_info) >= 2:
                    text = line_info[1][0]
                    lines.append(text)
        return "\n".join(lines)
    except ImportError:
        logger.warning("PaddleOCR not installed — image loading skipped")
        return ""
    except Exception as e:
        logger.error(f"OCR failed for {file_path}: {e}")
        return ""


def _load_code(file_path: Path) -> str:
    """Load source code with language tag."""
    ext = file_path.suffix
    lang_map = {
        ".py": "python", ".java": "java", ".go": "go",
        ".js": "javascript", ".ts": "typescript",
        ".cpp": "cpp", ".c": "c",
    }
    lang = lang_map.get(ext, ext.lstrip("."))
    content = file_path.read_text(encoding="utf-8", errors="replace")
    return f"```{lang}\n{content}\n```"


# ── Loader dispatch table ───────────────────────────────────────────

_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".txt": _load_txt,
    ".md": _load_markdown,
    ".html": _load_html,
    ".htm": _load_html,
    ".csv": _load_csv,
    ".xlsx": _load_excel,
    ".pptx": _load_pptx,
    ".json": _load_json,
    ".epub": _load_epub,
    ".png": _load_image_ocr,
    ".jpg": _load_image_ocr,
    ".jpeg": _load_image_ocr,
    ".py": _load_code,
    ".java": _load_code,
    ".go": _load_code,
    ".js": _load_code,
    ".ts": _load_code,
    ".cpp": _load_code,
    ".c": _load_code,
}


def load_file(file_path: str | Path) -> Document:
    """
    Load a single file and return a LlamaIndex Document.

    Raises ValueError for unsupported file types.
    """
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    if ext not in _LOADERS:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {sorted(_LOADERS.keys())}"
        )

    logger.info(f"Loading {ext}: {path.name}")
    text = _LOADERS[ext](path)

    Document = _get_Document()
    return Document(
        text=text,
        metadata={
            "file_name": path.name,
            "file_type": ext,
            "file_path": str(path),
        },
    )


def load_directory(directory: str | Path) -> list[Document]:
    """
    Load all supported files from a directory recursively.

    Skips unsupported extensions and files exceeding KB_MAX_FILE_SIZE_MB.
    """
    root = Path(directory).resolve()
    if not root.is_dir():
        raise NotADirectoryError(f"Not a directory: {root}")

    max_bytes = settings.KB_MAX_FILE_SIZE_MB * 1024 * 1024
    documents: list[Document] = []

    for path in root.rglob("*"):
        if not path.is_file():
            continue

        ext = path.suffix.lower()
        if ext not in _LOADERS:
            continue

        if path.stat().st_size > max_bytes:
            logger.warning(f"Skipping large file: {path.name} ({path.stat().st_size} bytes)")
            continue

        try:
            doc = load_file(path)
            documents.append(doc)
        except Exception as e:
            logger.error(f"Failed to load {path}: {e}")

    logger.info(f"Loaded {len(documents)} documents from {root}")
    return documents
